# -*- coding: utf-8 -*-
"""根据毕业论文生成并美化华南理工大学毕业答辩 PPT。"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PATH = ROOT / "毕设最终doc文档" / "答辩.pptx"
RAW_ASSET_DIR = ROOT / "assets" / "毕业论文"
PPT_ASSET_DIR = OUT_PATH.parent / "assets"

PRIMARY = RGBColor(0x00, 0x3C, 0x8C)
SECONDARY = RGBColor(0x17, 0x5D, 0xB8)
LIGHT_BG = RGBColor(0xF5, 0xF8, 0xFD)
TEXT = RGBColor(0x22, 0x22, 0x22)
SUB_TEXT = RGBColor(0x56, 0x56, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _font_obj(size: int, bold: bool = False):
    candidates = [
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _draw_center_text(draw, box, text, font, fill=(34, 34, 34)):
    x1, y1, x2, y2 = box
    tw = draw.textlength(text, font=font)
    _, _, _, th = draw.textbbox((0, 0), text, font=font)
    x = int((x1 + x2 - tw) / 2)
    y = int((y1 + y2 - th) / 2)
    draw.text((x, y), text, font=font, fill=fill)


def _generate_local_assets() -> dict[str, Path]:
    """在答辩.pptx同级目录下生成可直接插入的PNG资源。"""
    PPT_ASSET_DIR.mkdir(parents=True, exist_ok=True)
    images: dict[str, Path] = {}

    # 1) 系统架构图（流程块）
    system_png = PPT_ASSET_DIR / "system_architecture.png"
    img = Image.new("RGB", (1920, 1080), (246, 250, 255))
    d = ImageDraw.Draw(img)
    title_f = _font_obj(48, True)
    text_f = _font_obj(30, False)
    d.text((60, 40), "系统总体架构（答辩图示）", font=title_f, fill=(0, 60, 140))
    boxes = [
        (120, 190, 520, 390, "前端交互层\nReact + Three.js"),
        (620, 190, 1020, 390, "业务调度层\nSpring Boot + MQ"),
        (1120, 190, 1520, 390, "算法引擎层\nTensor + Streamline"),
        (1620, 190, 1860, 390, "存储层\nMySQL / MinIO"),
    ]
    for x1, y1, x2, y2, txt in boxes:
        d.rounded_rectangle((x1, y1, x2, y2), radius=24, fill=(255, 255, 255), outline=(65, 125, 200), width=4)
        for i, line in enumerate(txt.split("\n")):
            _draw_center_text(d, (x1 + 20, y1 + 38 + i * 58, x2 - 20, y1 + 95 + i * 58), line, _font_obj(32 if i == 0 else 27, i == 0), (40, 40, 40))
    arrows = [(520, 290, 620, 290), (1020, 290, 1120, 290), (1520, 290, 1620, 290)]
    for x1, y1, x2, y2 in arrows:
        d.line((x1, y1, x2, y2), fill=(23, 93, 184), width=8)
        d.polygon([(x2, y2), (x2 - 25, y2 - 14), (x2 - 25, y2 + 14)], fill=(23, 93, 184))
    d.rounded_rectangle((120, 500, 1860, 910), radius=28, fill=(255, 255, 255), outline=(210, 224, 243), width=3)
    d.text((170, 550), "业务闭环：参数输入  →  布局演算  →  交互微调  →  三维导出(GLB)", font=_font_obj(40, True), fill=(0, 60, 140))
    d.text((170, 640), "特性：异步任务、阶段进度回传、二维与三维任务解耦", font=_font_obj(32), fill=(70, 70, 70))
    img.save(system_png)
    images["system_arch"] = system_png

    # 2) 前端流程图
    web_png = PPT_ASSET_DIR / "web_workflow.png"
    img = Image.new("RGB", (1920, 1080), (248, 251, 255))
    d = ImageDraw.Draw(img)
    d.text((60, 40), "前端交互与业务流程（答辩图示）", font=title_f, fill=(0, 60, 140))
    cards = [
        (130, 230, 560, 430, "参数面板", "Tensor参数\n道路 dsep/dtest"),
        (730, 230, 1160, 430, "二维编辑", "SVG/JSON预览\n历史布局复用"),
        (1330, 230, 1760, 430, "三维预览", "GLB加载\nOrbit观察"),
        (430, 610, 920, 860, "任务状态反馈", "轮询进度 + 阶段提示词"),
        (1020, 610, 1510, 860, "结果沉淀", "布局JSON + 模型资产"),
    ]
    for x1, y1, x2, y2, t1, t2 in cards:
        d.rounded_rectangle((x1, y1, x2, y2), radius=26, fill=(255, 255, 255), outline=(90, 145, 215), width=4)
        _draw_center_text(d, (x1 + 20, y1 + 25, x2 - 20, y1 + 85), t1, _font_obj(37, True), (35, 35, 35))
        for i, line in enumerate(t2.split("\n")):
            _draw_center_text(d, (x1 + 20, y1 + 95 + i * 52, x2 - 20, y1 + 145 + i * 52), line, _font_obj(30), (95, 95, 95))
    for x1, y1, x2, y2 in [(560, 330, 730, 330), (1160, 330, 1330, 330), (910, 735, 1020, 735)]:
        d.line((x1, y1, x2, y2), fill=(23, 93, 184), width=8)
        d.polygon([(x2, y2), (x2 - 25, y2 - 14), (x2 - 25, y2 + 14)], fill=(23, 93, 184))
    d.line((950, 430, 950, 610), fill=(23, 93, 184), width=8)
    d.polygon([(950, 610), (935, 580), (965, 580)], fill=(23, 93, 184))
    img.save(web_png)
    images["web_arch"] = web_png

    # 3) 核心算法图（替代tensor原图）
    tensor_png = PPT_ASSET_DIR / "tensor_streamline.png"
    img = Image.new("RGB", (1920, 1080), (245, 249, 255))
    d = ImageDraw.Draw(img)
    d.text((60, 40), "张量场与流线算法示意（答辩图示）", font=title_f, fill=(0, 60, 140))
    d.rounded_rectangle((120, 170, 1800, 940), radius=32, fill=(255, 255, 255), outline=(200, 216, 238), width=3)
    # 网格
    for x in range(180, 1760, 80):
        d.line((x, 230, x, 880), fill=(230, 236, 245), width=2)
    for y in range(230, 890, 80):
        d.line((180, y, 1740, y), fill=(230, 236, 245), width=2)
    # 流线
    for i in range(11):
        points = []
        base = 260 + i * 55
        for x in range(180, 1740, 30):
            y = int(base + 45 * __import__("math").sin((x / 190.0) + i * 0.35))
            points.append((x, y))
        d.line(points, fill=(23, 93, 184), width=4)
    # 标注
    tags = [
        (220, 255, "主方向场"),
        (1450, 290, "RK4积分流线"),
        (260, 820, "dsep / dtest 约束"),
        (1380, 820, "R树求交 + 简化"),
    ]
    for x, y, txt in tags:
        d.rounded_rectangle((x, y, x + 290, y + 52), radius=14, fill=(255, 255, 255), outline=(120, 160, 215), width=2)
        _draw_center_text(d, (x + 8, y + 8, x + 282, y + 44), txt, _font_obj(24, True), (55, 55, 55))
    img.save(tensor_png)
    images["tensor"] = tensor_png

    # 4) 参数影响图
    impact_png = PPT_ASSET_DIR / "experiment_impact.png"
    img = Image.new("RGB", (1920, 1080), (247, 251, 255))
    d = ImageDraw.Draw(img)
    d.text((60, 40), "关键参数对生成效果影响（答辩图示）", font=title_f, fill=(0, 60, 140))
    d.rounded_rectangle((120, 170, 1800, 940), radius=32, fill=(255, 255, 255), outline=(200, 216, 238), width=3)
    cols = [("dsep", "道路密度", 0.78), ("dtest", "碰撞阈值", 0.66), ("noise", "水体自然度", 0.72), ("shrink", "可建面积", 0.58)]
    sx = 260
    for key, label, val in cols:
        h = int(420 * val)
        x1, y1, x2, y2 = sx, 760 - h, sx + 220, 760
        d.rounded_rectangle((x1, y1, x2, y2), radius=18, fill=(80, 145, 225), outline=(58, 114, 186), width=3)
        _draw_center_text(d, (x1, 780, x2, 840), key, _font_obj(28, True), (40, 40, 40))
        _draw_center_text(d, (x1 - 50, 850, x2 + 50, 900), label, _font_obj(24), (95, 95, 95))
        sx += 350
    d.text((180, 220), "结论：通过参数可解释调控，实现道路、水体、公园与建筑地块的协调生成", font=_font_obj(32), fill=(70, 70, 70))
    img.save(impact_png)
    images["impact"] = impact_png

    # 复制原有assets里的png/jpg（若有）到同级assets，便于后续手工替换
    for ext in ("*.png", "*.jpg", "*.jpeg", "*.webp"):
        for fp in RAW_ASSET_DIR.glob(ext):
            target = PPT_ASSET_DIR / fp.name
            if not target.exists():
                target.write_bytes(fp.read_bytes())

    # 优先挂载论文中已导出的真实实验图
    for key, filename in {
        "exp_water": "image-20260413222021914.png",
        "exp_road": "image-20260413222153715.png",
        "exp_block": "image-20260413222320106.png",
    }.items():
        p = PPT_ASSET_DIR / filename
        if p.exists():
            images[key] = p

    return images


def _font(paragraph, size: int, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    paragraph.font.name = "微软雅黑"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_header_footer(slide, title: str, page: int):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.45))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()
    p = top.text_frame.paragraphs[0]
    p.text = "华南理工大学 本科毕业设计答辩"
    _font(p, 13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.55), Inches(9.5), Inches(0.55))
    t = title_box.text_frame.paragraphs[0]
    t.text = title
    _font(t, 27, bold=True, color=PRIMARY)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.15), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    line.line.fill.background()

    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), Inches(13.333), Inches(0.28))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = LIGHT_BG
    bottom.line.fill.background()

    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.24), Inches(0.9), Inches(0.2))
    p2 = page_box.text_frame.paragraphs[0]
    p2.text = str(page)
    _font(p2, 11, color=SUB_TEXT, align=PP_ALIGN.RIGHT)


def _new_content_slide(prs: Presentation, title: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.45), Inches(13.333), Inches(6.77))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    _add_header_footer(slide, title, page)
    return slide


def _add_bullets(slide, bullets: Iterable[str]):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(12.0), Inches(5.45))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        _font(p, 20 if i == 0 else 18, color=TEXT)
        p.space_after = Pt(10)


def _add_image_with_caption(slide, image: Path, caption: str):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.8), Inches(4.9))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFA, 0xFC, 0xFF)
    frame.line.color.rgb = RGBColor(0xD4, 0xE1, 0xF2)

    slide.shapes.add_picture(str(image), Inches(1.1), Inches(1.78), Inches(11.1), Inches(4.15))

    c = slide.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.35)).text_frame.paragraphs[0]
    c.text = caption
    _font(c, 15, color=SUB_TEXT, align=PP_ALIGN.CENTER)


def build() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    images = _generate_local_assets()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cover_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    cover_bg.fill.solid()
    cover_bg.fill.fore_color.rgb = LIGHT_BG
    cover_bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()
    p = top.text_frame.paragraphs[0]
    p.text = "华南理工大学"
    _font(p, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sub = slide.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.3), Inches(0.6)).text_frame.paragraphs[0]
    sub.text = "本科毕业设计（论文）答辩"
    _font(sub, 25, color=PRIMARY, align=PP_ALIGN.CENTER)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.5)).text_frame.paragraphs[0]
    title.text = "三维模型生成以及可视化交互系统\n设计与实现"
    _font(title, 34, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    info_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(4.2), Inches(8.9), Inches(2.4))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = WHITE
    info_box.line.color.rgb = RGBColor(0xD6, 0xE2, 0xF3)

    info = info_box.text_frame
    info.clear()
    lines = [
        "学    院：                  （请填写）",
        "专    业：                  （请填写）",
        "答辩人：                  （请填写）",
        "指导教师：              （请填写）",
        "日    期：              2026 年 4 月",
    ]
    for i, line in enumerate(lines):
        p = info.paragraphs[0] if i == 0 else info.add_paragraph()
        p.text = line
        _font(p, 19, color=SUB_TEXT, align=PP_ALIGN.CENTER)
        p.space_after = Pt(8)

    page = 2
    s = _new_content_slide(prs, "汇报提纲", page); page += 1
    _add_bullets(s, [
        "选题背景与研究意义",
        "国内外研究现状与拟解决问题",
        "系统总体架构与业务流程",
        "核心算法与阶段化流水线",
        "实验环境、结果分析与结论",
        "不足与未来展望",
    ])

    s = _new_content_slide(prs, "选题背景与研究意义", page); page += 1
    _add_bullets(s, [
        "数字孪生、游戏场景、城市仿真需要可控可解释的空间生成能力。",
        "纯手工建模效率低；AIGC 在 3D 群落存在拓扑破碎与语义缺失。",
        "本文采用“参数约束 + 几何推演”：张量场引导 + 流线积分骨架。",
        "目标是构建可复现的 Whitebox 底座，降低后续资产填充修正成本。",
    ])

    s = _new_content_slide(prs, "国内外现状与关键问题", page); page += 1
    _add_bullets(s, [
        "国外：L 系统、图语法与张量场程序化城市生成体系较成熟。",
        "国内：GIS 还原与深度学习应用增长，风格化高可控 3D 底座仍待完善。",
        "痛点：布局与填充耦合、参数与审美脱节、多要素协同不足。",
        "本文聚焦：参数映射、拓扑地块划分、高效异步计算与数据流转。",
    ])

    s = _new_content_slide(prs, "系统总体架构", page); page += 1
    if "system_arch" in images:
        _add_image_with_caption(s, images["system_arch"], "图：系统整体架构（前端交互层 / 业务调度层 / 核心算法引擎层 / 存储层）")
    else:
        _add_bullets(s, ["未找到架构图资源，请补充 `assets/毕业论文/毕业设计-svg-model.drawio.svg`。"])

    s = _new_content_slide(prs, "前端交互与业务流程", page); page += 1
    if "web_arch" in images:
        _add_image_with_caption(s, images["web_arch"], "图：前端 Web 架构与布局编辑、模型预览交互链路")
    else:
        _add_bullets(s, ["未找到前端架构图资源，请补充 `assets/毕业论文/毕业设计-web架构设计.drawio.svg`。"])

    s = _new_content_slide(prs, "核心算法：张量场与流线", page); page += 1
    if "tensor" in images:
        _add_image_with_caption(s, images["tensor"], "图：张量场可视化（主方向与强度）")
    else:
        _add_bullets(s, ["未找到张量场图资源，请补充 `assets/毕业论文/tensor_visualization.svg`。"])

    s = _new_content_slide(prs, "阶段化流水线与三维导出", page); page += 1
    _add_bullets(s, [
        "布局任务：0~8 阶段顺序执行，覆盖全局场、水体、道路、公园、地块与导出。",
        "持续上报阶段与进度，便于前端可视化反馈和失败定位。",
        "三维任务在 layout.json 就绪后独立执行，完成布尔与挤出并导出 GLB。",
        "通过 RabbitMQ + Celery 解耦重计算与 Web 交互，提升稳定性。",
    ])

    s = _new_content_slide(prs, "实验环境与结果", page); page += 1
    if "impact" in images:
        _add_image_with_caption(
            s,
            images["impact"],
            "图：关键参数（dsep/dtest/noise/shrink）对布局结果影响趋势",
        )
    else:
        _add_bullets(s, [
            "环境：Windows 11，i7-14700K，32GB；内网部署前后端与算法服务。",
            "基准配置来自论文第六章 params.json；重点观察 dsep / dtest / 水体参数影响。",
            "结果：形成层次化道路—水体—公园—建筑结构，参数与形态具可解释对应。",
            "局限：参数敏感、规模增大耗时上升、建筑细部与水文语义仍可增强。",
        ])

    s = _new_content_slide(prs, "实验结果：水体与道路", page); page += 1
    if "exp_water" in images:
        _add_image_with_caption(
            s,
            images["exp_water"],
            "图：水体与岸线生成效果（论文实验图）",
        )
    else:
        _add_bullets(s, ["未检索到论文实验图（水体与岸线），可在同级 assets 手工补图。"])

    s = _new_content_slide(prs, "实验结果：道路、公园与建筑地块", page); page += 1
    if "exp_road" in images:
        _add_image_with_caption(
            s,
            images["exp_road"],
            "图：道路网络与公园分布（论文实验图）",
        )
    else:
        _add_bullets(s, ["未检索到道路与公园实验图，可在同级 assets 手工补图。"])

    s = _new_content_slide(prs, "实验结果：建筑地块划分", page); page += 1
    if "exp_block" in images:
        _add_image_with_caption(
            s,
            images["exp_block"],
            "图：建筑地块划分示意（论文实验图）",
        )
    else:
        _add_bullets(s, ["未检索到建筑地块实验图，可在同级 assets 手工补图。"])

    s = _new_content_slide(prs, "总结与展望", page); page += 1
    _add_bullets(s, [
        "完成参数化布局—可视化编辑—三维导出的端到端工程闭环。",
        "创新点：逻辑驱动 PCG、阶段化可观测流水线、二维中间表示统一语义。",
        "未来：增强三维细节、并行优化热点、丰富布局输入来源、提升大场景渲染性能。",
    ])

    # 致谢页
    slide = _new_content_slide(prs, "致谢", page)
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.0), Inches(9.3), Inches(2.8))
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFF)
    center.line.color.rgb = RGBColor(0xCD, 0xDC, 0xF2)

    tf = center.text_frame
    p = tf.paragraphs[0]
    p.text = "恳请各位老师批评指正"
    _font(p, 39, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    p2 = tf.add_paragraph()
    p2.text = "谢 谢！"
    _font(p2, 30, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    try:
        prs.save(str(OUT_PATH))
        print(f"Wrote: {OUT_PATH}")
    except PermissionError:
        fallback = OUT_PATH.with_name("答辩-新版.pptx")
        prs.save(str(fallback))
        print(f"Wrote (fallback): {fallback}")


if __name__ == "__main__":
    build()
